#!/usr/bin/env python3
import argparse
import os
import re
import shutil
import subprocess
import sys
from pathlib import Path
import time


def parse_pages(pages_str: str) -> list[int]:
    pages: set[int] = set()
    for token in pages_str.split(","):
        token = token.strip()
        if not token:
            continue
        if "-" in token:
            a, b = token.split("-", 1)
            try:
                start, end = int(a), int(b)
            except ValueError:
                raise SystemExit(f"Invalid page range: {token}")
            if start > end:
                start, end = end, start
            pages.update(range(start, end + 1))
        else:
            try:
                pages.add(int(token))
            except ValueError:
                raise SystemExit(f"Invalid page value: {token}")
    return sorted(pages)


def run_qlmanage(pptx_path: Path, tmp_dir: Path, size: int) -> None:
    cmd = [
        "qlmanage",
        "-t",  # generate thumbnails
        "-s",
        str(size),
        "-o",
        str(tmp_dir),
        str(pptx_path),
    ]
    try:
        proc = subprocess.run(
            cmd,
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
        )
    except FileNotFoundError:
        raise SystemExit("qlmanage not found. This script requires macOS Quick Look.")
    except subprocess.CalledProcessError as e:
        print(e.stdout or "", file=sys.stderr)
        raise SystemExit(f"qlmanage failed with exit code {e.returncode}")


def parse_generated_pages(tmp_dir: Path) -> dict[int, Path]:
    # Map page_number -> file path by extracting trailing integer from filenames
    mapping: dict[int, Path] = {}
    for p in tmp_dir.glob("*.png"):
        m = re.search(r"(\d+)(?=\D*$)", p.name)
        if not m:
            # Skip files that don't contain a page number
            continue
        page_num = int(m.group(1))
        mapping[page_num] = p
    if not mapping:
        # qlmanage on PPTX often yields a single file thumbnail like 'filename.pptx.png'.
        # Handle that as page 1 if present, otherwise report empty.
        single = list(tmp_dir.glob("*.png"))
        if len(single) == 1:
            return {1: single[0]}
        raise SystemExit(f"No PNGs were generated in {tmp_dir}")
    return dict(sorted(mapping.items()))


def _run_osascript(script: str) -> None:
    try:
        proc = subprocess.run(
            ["osascript", "-e", script],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
        )
    except FileNotFoundError:
        raise SystemExit("osascript not found. AppleScript is required for WPS automation.")
    except subprocess.CalledProcessError as e:
        print(e.stdout or "", file=sys.stderr)
        raise SystemExit(f"AppleScript execution failed with exit code {e.returncode}")


def export_with_wps_slideshow(pptx_path: Path, export_dir: Path, pages: list[int]) -> dict[int, Path]:
    """Use WPS Office slideshow + screen capture to export specified slides as PNG.

    - Launches WPS Office with the PPTX
    - Starts slideshow from beginning (F5)
    - Navigates to requested slides using Right Arrow
    - Captures full screen for each slide via `screencapture`
    """
    export_dir.mkdir(parents=True, exist_ok=True)

    # Try to open with WPS Office app bundle
    wps_app = "/Applications/wpsoffice.app"
    if not Path(wps_app).exists():
        raise SystemExit("WPS Office is not installed at /Applications/wpsoffice.app")
    try:
        subprocess.run(["open", "-a", wps_app, str(pptx_path)], check=True)
    except subprocess.CalledProcessError as e:
        raise SystemExit("Failed to open PPTX with WPS Office")

    # Wait for WPS process and focus it
    wait_script = (
        'tell application id "com.kingsoft.wpsoffice.mac" to activate\n'
        'tell application "System Events"\n'
        '  set wpsProc to missing value\n'
        '  repeat 40 times\n'
        '    try\n'
        '      set wpsProc to first process whose bundle identifier is "com.kingsoft.wpsoffice.mac"\n'
        '      exit repeat\n'
        '    end try\n'
        '    delay 0.25\n'
        '  end repeat\n'
        '  if wpsProc is missing value then error "WPS process not found"\n'
        '  set frontmost of wpsProc to true\n'
        'end tell'
    )
    _run_osascript(wait_script)
    time.sleep(0.8)

    # Start slideshow from beginning: F5 (key code 96)
    _run_osascript('tell application "System Events" to key code 96')
    time.sleep(1.0)

    page_to_file: dict[int, Path] = {}
    current = 1
    pages_sorted = sorted(pages)
    for p in pages_sorted:
        if p < 1:
            continue
        if p < current:
            # Restart show to go back to beginning
            _run_osascript('tell application "System Events" to key code 53')  # Esc
            time.sleep(0.5)
            _run_osascript('tell application "System Events" to key code 96')  # F5
            time.sleep(0.8)
            current = 1
        # Move forward delta slides
        delta = p - current
        if delta > 0:
            nav_script = (
                'tell application "System Events"\n'
                f'  repeat {delta} times\n'
                '    key code 124\n'  # Right Arrow
                '    delay 0.1\n'
                '  end repeat\n'
                'end tell'
            )
            _run_osascript(nav_script)
            current = p
            time.sleep(0.2)

        # Capture entire main display
        out_path = export_dir / f"p{p}.png"
        try:
            subprocess.run(["screencapture", "-x", str(out_path)], check=True)
        except FileNotFoundError:
            raise SystemExit("screencapture not found. macOS is required for screen capture.")
        except subprocess.CalledProcessError as e:
            raise SystemExit(f"screencapture failed with exit code {e.returncode}")
        page_to_file[p] = out_path
        time.sleep(0.2)

    # Exit slideshow
    _run_osascript('tell application "System Events" to key code 53')
    time.sleep(0.2)

    return page_to_file


def export_with_keynote(pptx_path: Path, export_dir: Path) -> None:
    """Use Apple Keynote to export slides as images via AppleScript.

    This does not modify the PPTX; Keynote opens and exports slide images.
    """
    export_dir.mkdir(parents=True, exist_ok=True)
    # AppleScript to open PPTX and export as slide images (PNG)
    # Avoid quoting issues by passing as a single script string.
    ascript = f'''
set theInput to POSIX file "{pptx_path}"
set theOutput to POSIX file "{export_dir}"
tell application "Keynote"
    activate
    set theDoc to open theInput
    delay 1
    try
        export theDoc to theOutput as slide images
    on error errMsg number errNum
        close theDoc saving no
        error errMsg number errNum
    end try
    close theDoc saving no
end tell
'''
    try:
        # Write AppleScript to a temporary file to avoid quoting issues
        script_path = export_dir / "_export_keynote.applescript"
        script_path.write_text(ascript, encoding="utf-8")
        proc = subprocess.run(
            ["osascript", str(script_path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
        )
    except FileNotFoundError:
        raise SystemExit("osascript not found. AppleScript is required for Keynote export.")
    except subprocess.CalledProcessError as e:
        print(e.stdout or "", file=sys.stderr)
        raise SystemExit(f"Keynote export failed with exit code {e.returncode}")


def export_pdf_with_powerpoint(pptx_path: Path, pdf_path: Path) -> None:
    """Use Microsoft PowerPoint (AppleScript) to export a PPTX to PDF."""
    pdf_path.parent.mkdir(parents=True, exist_ok=True)
    ascript = f'''
set theInput to (POSIX file "{pptx_path}") as alias
set thePDF to POSIX file "{pdf_path}"
tell application "Microsoft PowerPoint"
    activate
    open theInput
    set thePres to active presentation
    try
        save thePres in thePDF as save as PDF
    on error errMsg number errNum
        close thePres saving no
        error errMsg number errNum
    end try
    close thePres saving no
end tell
'''
    try:
        script_path = pdf_path.parent / "_export_powerpoint_pdf.applescript"
        script_path.write_text(ascript, encoding="utf-8")
        proc = subprocess.run(
            ["osascript", str(script_path)],
            check=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
        )
    except FileNotFoundError:
        raise SystemExit("osascript not found. AppleScript is required for PowerPoint export.")
    except subprocess.CalledProcessError as e:
        print(e.stdout or "", file=sys.stderr)
        raise SystemExit(f"PowerPoint PDF export failed with exit code {e.returncode}")


def main():
    parser = argparse.ArgumentParser(
        description=(
            "Export PPTX slides to PNG via macOS Quick Look without modifying the source file."
        )
    )
    parser.add_argument(
        "--input",
        default=str(Path("input") / "LRTBH.pptx"),
        help="Path to the source PPTX file",
    )
    parser.add_argument(
        "--output-dir",
        default=str(Path("images_direct")),
        help="Directory to place the exported images",
    )
    parser.add_argument(
        "--pages",
        default=None,
        help="Pages to export, e.g. '8' or '8,10-12'. Default: all",
    )
    parser.add_argument(
        "--size",
        type=int,
        default=2400,
        help="Max dimension in pixels for generated thumbnails (passed to qlmanage -s)",
    )

    args = parser.parse_args()

    pptx_path = Path(args.input).resolve()
    if not pptx_path.exists():
        raise SystemExit(f"Input PPTX not found: {pptx_path}")

    out_dir = Path(args.output_dir).resolve()
    out_dir.mkdir(parents=True, exist_ok=True)

    tmp_dir = out_dir / "_ql_tmp"
    if tmp_dir.exists():
        shutil.rmtree(tmp_dir)
    tmp_dir.mkdir(parents=True, exist_ok=True)

    # First attempt: Quick Look (qlmanage)
    run_qlmanage(pptx_path, tmp_dir, args.size)

    page_to_file = parse_generated_pages(tmp_dir)

    wanted = None
    if args.pages:
        wanted = set(parse_pages(args.pages))
        selected = [p for p in sorted(wanted) if p in page_to_file]
    else:
        selected = list(page_to_file.keys())

    # If Quick Look couldn't produce per-slide output for requested pages, fallback to Keynote
    need_fallback = False
    if args.pages:
        # missing any wanted pages
        if any(p not in page_to_file for p in wanted):
            need_fallback = True
    else:
        # exporting all slides but qlmanage gave only one image (thumbnail)
        if len(page_to_file) <= 1:
            need_fallback = True

    if need_fallback:
        # Attempt fallback 1: PowerPoint export to PDF, then qlmanage per-page PNGs
        shutil.rmtree(tmp_dir, ignore_errors=True)
        pp_pdf_dir = out_dir / "_pp_pdf"
        pdf_path = pp_pdf_dir / (pptx_path.stem + ".pdf")
        try:
            export_pdf_with_powerpoint(pptx_path, pdf_path)
            pdf_tmp = out_dir / "_pdf_tmp"
            if pdf_tmp.exists():
                shutil.rmtree(pdf_tmp)
            pdf_tmp.mkdir(parents=True, exist_ok=True)

            run_qlmanage(pdf_path, pdf_tmp, args.size)
            page_to_file = parse_generated_pages(pdf_tmp)
        except SystemExit as e:
            # Fallback 2: Keynote (if installed) export images directly
            kn_tmp = out_dir / "_keynote_tmp"
            if kn_tmp.exists():
                shutil.rmtree(kn_tmp)
            kn_tmp.mkdir(parents=True, exist_ok=True)
            try:
                export_with_keynote(pptx_path, kn_tmp)
                page_to_file = parse_generated_pages(kn_tmp)
            except SystemExit:
                # Fallback 3: WPS slideshow + screen capture (GUI automation)
                if not args.pages:
                    # We need target pages to drive slideshow navigation
                    raise
                try:
                    page_to_file = export_with_wps_slideshow(pptx_path, out_dir, sorted(wanted))
                except SystemExit:
                    # Fallback 4: python-pptx to build single-slide PPTX, then qlmanage
                    if not args.pages:
                        raise
                    try:
                        from pptx import Presentation  # type: ignore
                    except Exception:
                        raise SystemExit(
                            "python-pptx is required for fallback without AppleScript.\n"
                            "Install with: pip install python-pptx lxml"
                        )

                    page_to_file = {}
                    prs_tmp_dir = out_dir / "_prs_tmp"
                    if prs_tmp_dir.exists():
                        shutil.rmtree(prs_tmp_dir)
                    prs_tmp_dir.mkdir(parents=True, exist_ok=True)

                    wanted_pages = sorted(wanted)
                    for p in wanted_pages:
                        single_pptx = prs_tmp_dir / f"single_p{p}.pptx"
                        shutil.copy2(pptx_path, single_pptx)
                        prs = Presentation(str(single_pptx))
                        total = len(prs.slides)
                        idx_keep = p - 1
                        if idx_keep < 0 or idx_keep >= total:
                            continue
                        # delete other slides (from end to start)
                        for idx in reversed(range(total)):
                            if idx == idx_keep:
                                continue
                            rId = prs.slides._sldIdLst[idx].rId  # type: ignore[attr-defined]
                            prs.part.drop_rel(rId)
                            del prs.slides._sldIdLst[idx]  # type: ignore[attr-defined]
                        prs.save(str(single_pptx))

                        # export image for this single-slide pptx
                        single_tmp = prs_tmp_dir / f"_p{p}_tmp"
                        single_tmp.mkdir(parents=True, exist_ok=True)
                        run_qlmanage(single_pptx, single_tmp, args.size)
                        # Should create one png; map it to page p
                        pngs = list(single_tmp.glob("*.png"))
                        if pngs:
                            page_to_file[p] = pngs[0]

        if args.pages:
            selected = [p for p in sorted(wanted) if p in page_to_file]
        else:
            selected = list(page_to_file.keys())

    results: list[tuple[int, Path]] = []
    for p in selected:
        src = page_to_file[p]
        dst = out_dir / f"p{p}.png"
        try:
            if src.resolve() != dst.resolve():
                shutil.copy2(src, dst)
            else:
                # Source already at destination; ensure file exists
                if not dst.exists():
                    # If for some reason path is same but missing, write a copy
                    shutil.copy2(src, dst)
        except Exception:
            # As a safety, attempt to overwrite by moving temp file
            if src.exists() and src != dst:
                shutil.move(str(src), str(dst))
        results.append((p, dst))

    # Clean tmp
    shutil.rmtree(tmp_dir, ignore_errors=True)
    # cleanup any tmp folders created
    for d in [out_dir / "_keynote_tmp", out_dir / "_pdf_tmp", out_dir / "_pp_pdf", out_dir / "_prs_tmp"]:
        shutil.rmtree(d, ignore_errors=True)

    if not results:
        raise SystemExit("No slides exported.")

    print("Exported:")
    for p, d in results:
        print(f"  page {p} -> {d}")


if __name__ == "__main__":
    main()